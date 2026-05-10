from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / 'docs' / 'DeployFlow_Final_Presentation.pptx'

BG = RGBColor(7, 10, 18)
CARD = RGBColor(17, 24, 39)
LINE = RGBColor(55, 65, 81)
TEXT = RGBColor(248, 250, 252)
MUTED = RGBColor(164, 174, 192)
ACCENT = RGBColor(139, 92, 246)
CYAN = RGBColor(34, 211, 238)
GREEN = RGBColor(16, 185, 129)
RED = RGBColor(244, 63, 94)
YELLOW = RGBColor(245, 158, 11)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    return slide


def add_title(slide, title, subtitle=None, x=0.65, y=0.55, w=11.8, size=36):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.8))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = TEXT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(8)
    return box


def add_kicker(slide, text, x=0.7, y=0.38):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.2), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(8)
    p.font.bold = True
    p.font.color.rgb = CYAN
    return box


def card(slide, x, y, w, h, radius=True, fill=CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    return shape


def text(slide, content, x, y, w, h, size=16, color=TEXT, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align:
        p.alignment = align
    return box


def bullets(slide, items, x, y, w, h, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = MUTED
        p.space_after = Pt(8)
    return box


def metric(slide, label, value, x, y, color=CYAN):
    card(slide, x, y, 2.15, 1.28)
    text(slide, label.upper(), x+0.16, y+0.16, 1.8, 0.24, 8, MUTED, True)
    text(slide, value, x+0.16, y+0.44, 1.8, 0.45, 26, TEXT, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.16), Inches(y+1.05), Inches(0.65), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = color; line.line.fill.background()


def code(slide, content, x, y, w, h):
    card(slide, x, y, w, h, fill=RGBColor(10, 15, 29))
    box = slide.shapes.add_textbox(Inches(x+0.18), Inches(y+0.18), Inches(w-0.36), Inches(h-0.36))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.name = 'Consolas'
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(226, 232, 240)
    return box


def logo(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.66), Inches(6.82), Inches(0.42), Inches(0.32))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
    text(slide, 'DF', 0.72, 6.88, 0.3, 0.1, 8, TEXT, True, PP_ALIGN.CENTER)
    text(slide, 'DeployFlow', 1.14, 6.81, 2.0, 0.25, 9, MUTED, True)

# Slide 1
s = blank()
add_kicker(s, 'software deployment window planner')
add_title(s, 'DeployFlow', 'Graph coloring applied to daily release coordination', 0.7, 0.85, 7.0, 44)
text(s, 'Plan safe launch windows when services share teams, infrastructure, dependencies, and risk.', 0.72, 3.0, 6.1, 0.9, 20, MUTED)
metric(s, 'Selected services', '9', 8.1, 1.05, CYAN)
metric(s, 'Safe windows', '5', 10.45, 1.05, GREEN)
metric(s, 'Conflict edges', '15', 8.1, 2.65, YELLOW)
metric(s, 'Recursive calls', '23', 10.45, 2.65, ACCENT)
text(s, 'Artem Pasichnyk · Yaroslav Kondratenko', 0.72, 6.35, 7, 0.35, 12, MUTED)
logo(s)

# Slide 2
s = blank(); add_kicker(s, 'problem'); add_title(s, 'From map coloring to release planning')
card(s, 0.7, 1.55, 3.8, 4.8); card(s, 4.85, 1.55, 3.8, 4.8); card(s, 9.0, 1.55, 3.1, 4.8)
text(s, 'Map coloring', 1.0, 1.9, 3.0, 0.4, 20, TEXT, True)
bullets(s, ['Region = vertex', 'Shared border = edge', 'Color = label', 'Goal: adjacent regions differ'], 1.0, 2.5, 3.0, 2.5, 16)
text(s, 'DeployFlow', 5.15, 1.9, 3.0, 0.4, 20, TEXT, True)
bullets(s, ['Deployment = vertex', 'Unsafe parallel pair = edge', 'Window = color', 'Goal: conflicting releases differ'], 5.15, 2.5, 3.1, 2.5, 16)
text(s, 'Result', 9.3, 1.9, 2.5, 0.4, 20, TEXT, True)
bullets(s, ['Timeline', 'Runbook', 'Conflict graph', 'Search metrics'], 9.3, 2.5, 2.3, 2.5, 16)
logo(s)

# Slide 3
s = blank(); add_kicker(s, 'classical algorithm'); add_title(s, 'Backtracking graph coloring')
code(s, '''COLOR_VERTEX(i):\n  if i == n: return true\n\n  for color in 1..k:\n    if SAFE(i, color):\n      colors[i] = color\n      if COLOR_VERTEX(i + 1):\n        return true\n      colors[i] = 0\n\n  return false''', 0.8, 1.55, 5.8, 4.6)
card(s, 7.0, 1.55, 5.3, 4.6)
text(s, 'How it behaves', 7.35, 1.9, 4, 0.4, 22, TEXT, True)
bullets(s, ['Systematically explores color assignments.', 'Rejects a branch as soon as a neighbor conflict appears.', 'Complete for the chosen number of windows.', 'Can be expensive on dense graphs.'], 7.35, 2.55, 4.5, 2.8, 17)
logo(s)

# Slide 4
s = blank(); add_kicker(s, 'product'); add_title(s, 'A useful release board, not a toy demo')
img = ROOT / 'docs' / 'screenshots' / '01-dashboard.png'
s.shapes.add_picture(str(img), Inches(0.7), Inches(1.35), width=Inches(6.9))
card(s, 8.0, 1.35, 4.4, 5.1)
text(s, 'What the user does', 8.35, 1.7, 3.4, 0.4, 22, TEXT, True)
bullets(s, ['Select today’s service updates.', 'Add custom deployments.', 'Choose available windows.', 'Switch classic vs smart solver.', 'Generate a release timeline.'], 8.35, 2.3, 3.6, 3.1, 16)
logo(s)

# Slide 5
s = blank(); add_kicker(s, 'architecture'); add_title(s, 'One Java application, product-grade layers')
items = [
    ('Web UI', 'resources/web', 'dashboard, timeline, graph view'),
    ('HTTP driver', 'DeployFlowApp.java', 'serves UI and API'),
    ('Planner', 'DeploymentPlanner.java', 'builds conflict graph'),
    ('Algorithm', 'GraphColoringSolver.java', 'classic + improved solver'),
]
for idx, (name, path, desc) in enumerate(items):
    x = 0.8 + idx * 3.1
    card(s, x, 1.7, 2.7, 3.9)
    text(s, name, x+0.2, 2.0, 2.2, 0.35, 20, TEXT, True)
    text(s, path, x+0.2, 2.55, 2.2, 0.4, 12, CYAN, True)
    text(s, desc, x+0.2, 3.2, 2.15, 1.0, 15, MUTED)
    if idx < len(items)-1:
        text(s, '→', x+2.75, 3.2, 0.45, 0.5, 30, CYAN, True, PP_ALIGN.CENTER)
text(s, 'No Maven, Gradle, Node, or external dependencies are required to run the app.', 0.85, 6.1, 11, 0.4, 15, MUTED)
logo(s)

# Slide 6
s = blank(); add_kicker(s, 'improved algorithm'); add_title(s, 'Smarter search for dense deployment days')
card(s, 0.8, 1.45, 5.5, 4.9)
text(s, 'Classical', 1.15, 1.8, 2.2, 0.4, 22, TEXT, True)
bullets(s, ['Fixed vertex order', 'Try colors 1..k', 'Only checks colored neighbors', 'May explore weak branches'], 1.15, 2.4, 4.4, 2.7, 17)
card(s, 6.75, 1.45, 5.5, 4.9)
text(s, 'DeployFlow Smart', 7.1, 1.8, 3.4, 0.4, 22, TEXT, True)
bullets(s, ['Most constrained vertex first', 'Degree + risk tie-breakers', 'Least-constraining color', 'Forward checking', 'Dependency windows must be chronological'], 7.1, 2.4, 4.6, 3.0, 17)
logo(s)

# Slide 7
s = blank(); add_kicker(s, 'demo result'); add_title(s, 'Generated release timeline')
img = ROOT / 'docs' / 'screenshots' / '02-release-plan-report.png'
s.shapes.add_picture(str(img), Inches(0.7), Inches(1.25), width=Inches(6.0))
card(s, 7.35, 1.25, 5.0, 5.3)
text(s, 'Output views', 7.7, 1.6, 3.5, 0.4, 22, TEXT, True)
bullets(s, ['Release windows with parallel lanes.', 'Conflict graph colored by assigned window.', 'Operator runbook with step order.', 'Metrics: calls, checks, backtracks, runtime.'], 7.7, 2.25, 4.1, 2.7, 17)
metric(s, 'Windows used', '5/6', 7.7, 5.0, GREEN)
metric(s, 'Runtime', '10.82ms', 10.05, 5.0, CYAN)
logo(s)

# Slide 8
s = blank(); add_kicker(s, 'execution'); add_title(s, 'Correctness checks visible in the product')
card(s, 0.8, 1.45, 4.0, 4.9)
text(s, 'Conflict rules', 1.15, 1.8, 3.0, 0.4, 22, TEXT, True)
bullets(s, ['Same owner capacity', 'Shared resource', 'Direct dependency', 'High-risk pair'], 1.15, 2.4, 3.2, 2.8, 17)
card(s, 5.0, 1.45, 3.9, 4.9)
text(s, 'Example metrics', 5.35, 1.8, 3.0, 0.4, 22, TEXT, True)
bullets(s, ['9 vertices', '15 conflict edges', '9 precedence rules', '23 recursive calls', '24 backtracks'], 5.35, 2.4, 3.0, 2.8, 17)
card(s, 9.1, 1.45, 3.4, 4.9)
text(s, 'Result', 9.45, 1.8, 2.5, 0.4, 22, TEXT, True)
bullets(s, ['Feasible plan', 'Chronological dependencies', 'No adjacent same window', 'Trace available'], 9.45, 2.4, 2.5, 2.8, 17)
logo(s)

# Slide 9
s = blank(); add_kicker(s, 'risks'); add_title(s, 'Risks, limits, and mitigations')
rows = [
    ('Dense graph', 'Backtracking can grow exponentially', 'MRV, degree ordering, forward checking'),
    ('Bad input data', 'Wrong edges create unsafe plans', 'Show conflict reasons and editable catalog'),
    ('Too few windows', 'No feasible coloring exists', 'Report failure and suggest extra windows'),
    ('Real-world constraints', 'People and maintenance windows vary', 'Add more domain rules as edges/constraints'),
]
y = 1.45
for risk, why, fix in rows:
    card(s, 0.8, y, 11.7, 1.1)
    text(s, risk, 1.1, y+0.2, 2.2, 0.35, 18, TEXT, True)
    text(s, why, 3.6, y+0.2, 3.6, 0.5, 14, MUTED)
    text(s, fix, 7.7, y+0.2, 4.1, 0.5, 14, CYAN)
    y += 1.25
logo(s)

# Slide 10
s = blank(); add_kicker(s, 'final'); add_title(s, 'Final thoughts')
card(s, 0.9, 1.55, 5.7, 4.7)
text(s, 'What we built', 1.25, 1.9, 3, 0.4, 24, TEXT, True)
bullets(s, ['Working Java product', 'Web dashboard and API', 'Classical graph coloring function', 'Improved planner algorithm', 'Screenshots, report, and demo-ready run scripts'], 1.25, 2.55, 4.6, 3.0, 17)
card(s, 7.0, 1.55, 5.1, 4.7)
text(s, 'Why it matters', 7.35, 1.9, 3, 0.4, 24, TEXT, True)
bullets(s, ['Shows map coloring as a real scheduling engine.', 'Makes deployment conflicts visible.', 'Gives teams a practical release order.', 'Leaves room for future production constraints.'], 7.35, 2.55, 4.0, 3.0, 17)
text(s, 'DeployFlow', 4.8, 6.55, 3.0, 0.4, 16, CYAN, True, PP_ALIGN.CENTER)
logo(s)

prs.save(OUT)
print(OUT)
